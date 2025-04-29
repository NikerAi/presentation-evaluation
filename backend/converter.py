from pdf2image import convert_from_bytes
from PIL import Image
from bs4 import BeautifulSoup
from pptx.enum.shapes import PP_PLACEHOLDER_TYPE
from pptx import Presentation
from pptx.shapes.base import _PlaceholderFormat
import subprocess
import base64
import io
import tempfile
import os
import zipfile


class GenImage():
    '''
    Class for working with a presentation in the form of an image

    Attributes
    ----------
    buffer: io.BytesIO
        buffer into which the image is loaded upon successful conversion
    default_fonts: Dict[str, str]
        Storage of standard fonts for the theme
    fonts: Dict[str, list]
    
    Methods
    -------
    pptx(pptx_bytes:bytes):
        Convert pptx to jpeg
    pdf(pdf_bytes: bytes):
        Convert pdf to jpeg
    parse_default_fonts(path_pptx):
        Parsing the standard fonts of the presentation theme
    parse_fonts_on_slide(path_pptx):
        Parsing custom fonts of slides
    base64() -> bytes:
        Get base64 bytes from buffer

    '''
    def __init__(self, bytes: bytes, file_format: str):
        '''
        Creates an image from the resulting byte array 
        '''
        self.buffer = io.BytesIO()
        self.default_fonts = {}
        self.fonts = {}
        # Automatically calling the converter, if the type is not supported, we throw an exception
        converter = getattr(self, file_format.lower(), lambda bytes: self.not_support(file_format))
        converter(bytes)

    def not_support(self, file_format: str):
        raise Exception(f"Not support this file format {file_format}")

    def parse_default_fonts(self, path_pptx):
        '''
        Parsing the standard fonts of the presentation theme

        Parameters
        ----------
            path_pptx: str
                Path to pptx file
        '''
        with zipfile.ZipFile(path_pptx) as z:
            theme_files = list(filter(lambda f: f.startswith("ppt/theme/") and f.endswith(".xml"), z.namelist()))
            if len(theme_files) > 0:
                with z.open(theme_files[0]) as theme:
                    soup = BeautifulSoup(theme.read(), 'xml')

                    major = soup.find('a:majorFont')
                    minor = soup.find('a:minorFont')

                    if major and major.find('a:latin'):
                        self.default_fonts['major'] = major.find('a:latin')['typeface']
                    if minor and minor.find('a:latin'):
                        self.default_fonts['minor'] = minor.find('a:latin')['typeface']


    # Called if there is no font in the placeholder.text
    def get_theme_font(self, type: PP_PLACEHOLDER_TYPE) -> str | None:
        '''
        Take the font from the default theme
        
        Parameters
        ----------
            type: PP_PLACEHOLDER_TYPE
                Placeholder location on the slide
        '''
        if 'minor' in self.default_fonts and type != PP_PLACEHOLDER_TYPE.TITLE:
            return self.default_fonts['minor']
        if 'major' in self.default_fonts and type == PP_PLACEHOLDER_TYPE.TITLE:
            return self.default_fonts['major']
        return None
    
    def parse_fonts_on_slide(self, path_pptx):
        '''
        Parsing custom fonts of slides

        Parameters
        ----------
            path_pptx: str
                Path to pptx file
        '''
        prs = Presentation(path_pptx)
        
        for slide_num, slide in enumerate(prs.slides, start=1):
            slide_fonts = []

            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue

                # Return Enum PlaceholderType
                # It is necessary to determine the position of the text on the slide
                placeholder_type = shape.placeholder_format.type if shape.is_placeholder else None  

                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        # If there is no custom font, we try to get it from the theme.
                        font_name = run.font.name if run.font.name is not None else self.get_theme_font(placeholder_type)

                        slide_fonts.append(font_name)

            if slide_fonts:
                self.fonts[f"Слайд {slide_num}"] = slide_fonts

    # It is automatically invoked if necessary
    def pptx(self, pptx_bytes):
        '''
        Convert pptx to jpeg

        Parameters
        ----------
            pptx_bytes: bytes
                An array of bytes that may contain a pptx
        '''
        with tempfile.TemporaryDirectory() as tmpdir:
            pptx_path = os.path.join(tmpdir, "presentation.pptx")
            pdf_path = os.path.join(tmpdir, "presentation.pdf")

            with open(pptx_path, "wb") as f:
                f.write(pptx_bytes)

            self.parse_default_fonts(pptx_path)
            self.parse_fonts_on_slide(pptx_path)

            # Run libreoffice for convert pptx to pdf
            subprocess.run([
                "libreoffice",
                "--headless",
                "--convert-to", "pdf",
                "--outdir", tmpdir,
                pptx_path
            ], check=True)

            os.remove(pptx_path)

            with open(pdf_path, "rb") as f:
                pdf_bytes = f.read()
            
            os.remove(pdf_path)

            # Run converter pdf to img
            self.pdf(pdf_bytes)

    # It is automatically invoked if necessary
    def pdf(self, pdf_bytes: bytes):
        '''
        Convert pdf to jpeg

        Parameters
        ----------
            pdf_bytes: bytes
                An array of bytes that may contain a pdf
        '''
        images = convert_from_bytes(pdf_bytes)
        sum_width = sum([img.width for img in images])
        max_height = max([img.height for img in images])
        # Create empty image for pasting
        self.img = Image.new("RGB", (sum_width, max_height), color=(255, 255, 255))
        x_offset = 0
        # Construction all the slides into one picture
        for img in images:
            self.img.paste(img, (x_offset, 0))
            x_offset += img.width
        self.img.save(self.buffer, "jpeg")

    def base64(self):
        '''
        Get base64 bytes from buffer
        '''
        return base64.b64encode(self.buffer.getvalue())
    
# For future changes towards safe conversion
def convert_to_img(file: bytes, format: str) -> GenImage:
    '''
    Function of converting a presentation into an image

    Parameters
    ----------
        file: bytes
            file to be converted, in bytes
        format: str
            format/type of file that needs to be converted
    '''
    return GenImage(file, format)