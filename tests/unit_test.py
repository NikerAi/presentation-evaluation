import pytest
from io import BytesIO
from backend.converter import GenImage
from PIL import Image
import base64
from pptx import Presentation
from pptx.util import Inches
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
import tempfile
import os
import asyncio
from streamlit.testing.v1 import AppTest
from backend.llm_call import send_request
from backend.converter import response_handler
from pathlib import Path

def create_simple_presentation(output_path):
    '''
    Create simple presentation for tests
    '''
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)
    slide.shapes.title.text = "Тестовая презентация"
    slide.placeholders[1].text = "Создано для теста GenImage"

    slide_layout = prs.slide_layouts[5]
    slide2 = prs.slides.add_slide(slide_layout)
    slide2.shapes.title.text = "Второй слайд"

    textbox = slide2.shapes.add_textbox(Inches(1), Inches(2), Inches(4), Inches(1))
    tf = textbox.text_frame
    p = tf.add_paragraph()
    run = p.add_run()
    run.text = "Текст со шрифтом Arial"
    run.font.name = "Arial"
    run.font.size = Pt(18)
    run.font.color.rgb = RGBColor(0, 0, 0)

    prs.save(output_path)

@pytest.fixture
def sample_pptx_bytes():
    '''
    Get bytes from simple pptx 
    '''
    with tempfile.NamedTemporaryFile(suffix=".pptx", delete=True) as tmp:
        create_simple_presentation(tmp.name)
        tmp.seek(0)
        data = tmp.read()
    return data


@pytest.fixture
def run_app():
    return AppTest.from_file("../frontend/app.py", default_timeout=30).run()

def test_pptx_conversion(sample_pptx_bytes):
    '''
    Checking that the buffer contains an image
    '''
    gi = GenImage(sample_pptx_bytes, "pptx")
    img = Image.open(gi.buffer)
    assert img.format == "JPEG"
    assert img.size[0] > 0 and img.size[1] > 0


def test_base64_output(sample_pptx_bytes):
    '''
    Checking that the base64 bytes contain a picture
    '''
    gi = GenImage(sample_pptx_bytes, "pptx")
    b64 = gi.base64()
    assert isinstance(b64, bytes)
    decoded = base64.b64decode(b64)
    # Check JPEG magic bytes
    assert decoded[:2] == b'\xff\xd8'

def test_unsupported_format():
    '''
    Checking the reaction to an unsupported format
    '''
    with pytest.raises(Exception) as exc:
        GenImage(BytesIO(b"dummy").getvalue(), "txt")
    assert "Not support this file format txt" in str(exc.value)

def test_fonts_parsing_from_sample_pptx(sample_pptx_bytes):
    gen = GenImage(sample_pptx_bytes, "pptx")

    # check default fonts (major/minor)
    assert isinstance(gen.default_fonts, dict)
    assert "major" in gen.default_fonts or "minor" in gen.default_fonts

    # Check parse fonts on slide
    assert isinstance(gen.fonts, dict)
    assert len(gen.fonts) >= 2

    # Checking for Arial in the fonts of the second slide
    all_fonts = [font for fonts in gen.fonts.values() for font in fonts]
    assert any("Arial" in font for font in all_fonts if font)

async def async_convert(sample_pptx_bytes):
    loop = asyncio.get_event_loop()
    return await loop.run_in_executor(None, GenImage, sample_pptx_bytes, "pptx")

@pytest.mark.asyncio
async def test_asyncio_parallel(sample_pptx_bytes):
    tasks = [async_convert(sample_pptx_bytes) for _ in range(3)]
    results = await asyncio.gather(*tasks)
    assert all(len(r.default_fonts) > 0 for r in results)


def test_show_prompt(run_app):
    """
    Checks if default prompt shown correctly
    """

    # at.tabs[1].text_area[0].input(prompt).run()
    run_app.tabs[1].button[0].click().run()
    assert run_app.session_state["prompt"][-10:] == "улучшения]"


def test_save_prompt(run_app):
    """
    Checks if user's prompt saved correctly, and empty prompt raises an error
    """
    # Add new prompt
    run_app.tabs[1].text_area[0].input("New prompt").run()
    run_app.tabs[1].button[0].click().run()
    assert run_app.session_state["prompt"] == "New prompt"
    # Check if prompt is empty
    run_app.tabs[1].text_area[0].input("     ").run()
    run_app.tabs[1].button[0].click().run()
    assert "Добавьте текст запроса, поле не может быть пустым" in run_app.error.values[0]

def test_send_request(sample_pptx_bytes):
    """
    Send a request to llm and check if response contains keywords from the new prompt
    """
    project_root = Path(__file__).parent.parent  
    prompt_path = project_root / "frontend" / "default_prompt.txt"
    
    with open(prompt_path, "r", encoding="utf-8") as f:
        prompt_utf8 = f.read()
    
    prompt = prompt_utf8.encode("windows-1251").decode("windows-1251")
    
    response = send_request(prompt=prompt, presentation=sample_pptx_bytes, file_format="pptx").choices[0].message.content
    
    assert response, "LLM returned an empty response"
    
    expected_keywords = ["Титульный слайд", "Введение", "Содержание", "Оформление", "Заключение и контакты", "Формат", "Общие рекомендации"]
    for keyword in expected_keywords:
        assert keyword in response, f"Response does not contain expected keyword '{keyword}'"

def test_response_handler():
    """
    Checks if function returns correct result
    """
    content = "Test message"
    text = response_handler(content)
    assert len(text) > 0
